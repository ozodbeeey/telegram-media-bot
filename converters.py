import os
import io
import fitz  # pymupdf
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches
import openpyxl
from docx import Document
from docx2pdf import convert as docx2pdf_convert
import comtypes.client


async def convert_pdf_to_word(input_path, output_path):
    """Converts PDF to Word (DOCX)."""
    cv = Converter(input_path)
    cv.convert(output_path, start=0, end=None)
    cv.close()
    return output_path

async def convert_pdf_to_pptx(input_path, output_path):
    """Converts PDF to PowerPoint (PPTX) by rendering pages as images."""
    doc = fitz.open(input_path)
    prs = Presentation()
    
    # Standard 4:3 slide size
    # width = Inches(10) 
    # height = Inches(7.5)
    
    # We can adjust slide size to match PDF page size if needed, 
    # but for simplicity let's fit the image to the slide.
    
    for page in doc:
        # Render page to image
        pix = page.get_pixmap(dpi=150)
        img_data = pix.tobytes("png")
        image_stream = io.BytesIO(img_data)
        
        slide_layout = prs.slide_layouts[6] # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add image to slide, filling it
        # Using default slide dimensions 10x7.5 inches
        slide.shapes.add_picture(image_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
    prs.save(output_path)
    doc.close()
    return output_path

async def convert_word_to_excel(input_path, output_path):
    """Extracts tables from Word and saves them to Excel."""
    doc = Document(input_path)
    wb = openpyxl.Workbook()
    
    # Remove default sheet if we have tables, or use it as first
    if not doc.tables:
        # No tables found
        ws = wb.active
        ws.title = "No Tables Found"
        ws['A1'] = "Ushbu Word faylda hech qanday jadval topilmadi."
    else:
        # Remove default sheet
        wb.remove(wb.active)
        
        for i, table in enumerate(doc.tables):
            ws = wb.create_sheet(title=f"Table {i+1}")
            rows = table.rows
            for r_idx, row in enumerate(rows):
                cells = row.cells
                for c_idx, cell in enumerate(cells):
                    ws.cell(row=r_idx+1, column=c_idx+1, value=cell.text.strip())
                    
    wb.save(output_path)
    return output_path

async def convert_excel_to_word(input_path, output_path):
    """Converts Excel data to Word tables."""
    wb = openpyxl.load_workbook(input_path)
    doc = Document()
    
    sheet = wb.active # Use active sheet
    
    doc.add_heading(f"Data from {os.path.basename(input_path)}", 0)
    
    # Calculate rows and columns
    rows = list(sheet.rows)
    if not rows:
        doc.add_paragraph("Excel fayl bo'sh.")
    else:
        num_rows = len(rows)
        num_cols = len(rows[0])
        
        table = doc.add_table(rows=num_rows, cols=num_cols)
        table.style = 'Table Grid'
        
        for r_idx, row in enumerate(rows):
            for c_idx, cell in enumerate(row):
                val = str(cell.value) if cell.value is not None else ""
                table.cell(r_idx, c_idx).text = val
                
    doc.save(output_path)
    return output_path
    doc.save(output_path)
    return output_path

async def convert_word_to_pdf(input_path, output_path):
    """Converts Word to PDF using docx2pdf (Requires MS Word)."""
    # docx2pdf is synchronous and uses COM, better to run in executor or just call it.
    # It handles opening Word instance.
    try:
        docx2pdf_convert(input_path, output_path)
        return output_path
    except Exception as e:
        raise Exception(f"Word to PDF Conversion Failed: {e}")

async def convert_pptx_to_pdf(input_path, output_path):
    """Converts PPTX to PDF using PowerPoint COM automation."""
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    
    try:
        deck = powerpoint.Presentations.Open(input_path)
        # 32 = ppSaveAsPDF
        deck.SaveAs(output_path, 32)
        deck.Close()
        return output_path
    except Exception as e:
        raise Exception(f"PPTX to PDF Failed: {e}")
    finally:
        # Don't quit app if existing, but usually safer to quit if we created it. 
        # Check if we should quit? Ideally yes to clean up.
        try: powerpoint.Quit()
        except: pass

async def convert_excel_to_pdf(input_path, output_path):
    """Converts Excel to PDF using Excel COM automation."""
    excel = comtypes.client.CreateObject("Excel.Application")
    excel.Visible = 0
    
    try:
        wb = excel.Workbooks.Open(input_path)
        # 0 = xlTypePDF
        wb.ExportAsFixedFormat(0, output_path)
        wb.Close(False)
        return output_path
    except Exception as e:
        raise Exception(f"Excel to PDF Failed: {e}")
    finally:
         try: excel.Quit()
         except: pass

async def convert_pptx_to_word(input_path, output_path):
    """Extracts text from PPTX slides to Word."""
    prs = Presentation(input_path)
    doc = Document()
    
    doc.add_heading(f"Text from {os.path.basename(input_path)}", 0)
    
    for i, slide in enumerate(prs.slides):
        doc.add_heading(f"Slide {i+1}", level=1)
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                doc.add_paragraph(shape.text)
                
    doc.save(output_path)
    return output_path

async def convert_word_to_pptx(input_path, output_path):
    """Converts Word paragraphs to PPTX slides."""
    doc = Document(input_path)
    prs = Presentation()
    
    # 6 is Empty layout, 1 is Title and Content
    layout = prs.slide_layouts[1] 
    
    for para in doc.paragraphs:
        if not para.text.strip():
            continue
            
        slide = prs.slides.add_slide(layout)
        # Set title only for now, or content
        if slide.shapes.title:
            slide.shapes.title.text = para.text[:50] + "..." if len(para.text) > 50 else para.text
        
        # If we want full text in body
        if len(slide.placeholders) > 1:
             slide.placeholders[1].text = para.text
             
    prs.save(output_path)
    return output_path
