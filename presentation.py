import wikipedia
import requests
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO

def generate_presentation(topic, output_path):
    """
    Generates a PowerPoint presentation from Wikipedia based on the topic.
    """
    # 1. Fetch content
    wikipedia.set_lang('uz')
    try:
        page = wikipedia.page(topic)
    except Exception:
        # Fallback to English
        wikipedia.set_lang('en')
        try:
            page = wikipedia.page(topic)
        except Exception as e:
            return False, f"Ma'lumot topilmadi: {e}"

    # 2. Create Presentation
    prs = Presentation()
    
    # --- SLIDE 1: Title ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = page.title
    subtitle.text = "Tayyorladi: Bot (Sizning ismingiz)"

    # --- SLIDE 2: Introduction ---
    # Bullet layout
    bullet_slide_layout = prs.slide_layouts[1] 
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    title_shape.text = "Kirish"
    
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = page.summary[:500] + "..." # Limit text

    # --- Try to add Image to Slide 2 ---
    if page.images:
        img_url = None
        valid_extensions = ('.jpg', '.jpeg', '.png')
        for img in page.images:
             if img.lower().endswith(valid_extensions):
                 img_url = img
                 break
        
        if img_url:
            try:
                response = requests.get(img_url)
                if response.status_code == 200:
                    image_stream = BytesIO(response.content)
                    # Add picture to right side
                    slide.shapes.add_picture(image_stream, Inches(6), Inches(2), width=Inches(3.5))
            except:
                pass

    # --- SLIDE 3+: Content Body ---
    # We'll take the content, split by headings, and make a few slides
    # Limited to 3-4 slides to be fast
    
    content = page.content.split('\n\n')
    slide_count = 0
    
    for section in content:
        if slide_count >= 4: break # Limit slides
        
        section = section.strip()
        if not section or len(section) < 50: continue
        
        # Check if it looks like a heading
        if section.startswith('==') and section.endswith('=='):
            # Start new slide
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            title_text = section.replace('=', '').strip()
            title_shape.text = title_text
            
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame
            slide_count += 1
            
        elif slide_count > 0:
            # Add text to current slide
            # Just take first 300 chars to avoid overflow
            try:
                # Get the last slide added
                curr_slide = prs.slides[-1]
                # Check if it has a body placeholder
                if len(curr_slide.placeholders) > 1:
                     body = curr_slide.placeholders[1]
                     if not body.text:
                         body.text = section[:400] + "..."
            except:
                pass

    # --- SLIDE LAST: Conclusion ---
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "E'tiboringiz uchun rahmat!"
    
    prs.save(output_path)
    return True, "Slayd tayyor."
